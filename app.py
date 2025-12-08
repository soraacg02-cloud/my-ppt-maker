import streamlit as st
import streamlit.components.v1 as components
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from io import BytesIO
import docx
from docx.document import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
import fitz  # PyMuPDF
import re
import pandas as pd

# --- 設定網頁標題 ---
st.set_page_config(page_title="PPT 重組生成器 (Prompt更新版)", page_icon="📑", layout="wide")
st.title("📑 PPT 重組生成器 (含獨立項 Claim 功能)")
st.caption("支援多檔上傳、自動排序、錯誤診斷，並可選擇是否產生「獨立項 Claim」分頁。")

# === NBLM 提示詞區塊 (更新：修改為最新的 6 點要求) ===
nblm_prompt = """根據上傳的所有來源，分開整理出以下重點(不要表格)：

1. 案號 / 日期 / 公司： *(案號依據"公開號"、日期依據"優先權日"、公司依據"申請人")
2. 解決問題：
3. 發明精神：*(不要有公式)
4. 一句重點： *(用來描述發明特徵重點，20字)
5. 代表圖：*(根據發明精神建議3張最可以說明發明精神的圖片，範例:FIG.3)
6. 獨立項claim： *(分組且分行條列式+對應的代表圖，claim要有位階縮排而且claim的元件要有標號)"""

st.info("💡 **NBLM 使用提示詞** (已更新為最新 6 點要求，點擊下方綠色按鈕一鍵複製)")

# 使用 HTML 建立顯眼複製按鈕
components.html(
    f"""
    <html>
    <head><meta charset="utf-8"></head>
    <body style="font-family: sans-serif; margin: 0; padding: 0;">
        <div style="display: flex; flex-direction: column; align-items: flex-start;">
            <textarea id="copyTarget" style="opacity: 0; position: absolute; z-index: -1;">{nblm_prompt}</textarea>
            <div style="background-color: #f0f2f6; padding: 15px; border-radius: 10px; white-space: pre-wrap; font-size: 14px; color: #31333F; border: 1px solid #d6d6d6; width: 95%; margin-bottom: 10px;">{nblm_prompt}</div>
            <button onclick="copyFunction()" style="background-color: #00CC66; color: white; border: none; padding: 12px 24px; font-size: 16px; font-weight: bold; border-radius: 8px; cursor: pointer; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">📋 點我一鍵複製提示詞</button>
            <span id="statusParams" style="color: #00CC66; font-weight: bold; margin-left: 10px; opacity: 0; transition: opacity 0.5s;">✅ 複製成功！</span>
        </div>
        <script>
        function copyFunction() {{
            var copyText = document.getElementById("copyTarget");
            copyText.select();
            navigator.clipboard.writeText(copyText.value).then(function() {{
                var status = document.getElementById("statusParams");
                status.style.opacity = '1';
                setTimeout(function(){{ status.style.opacity = '0'; }}, 2000);
            }});
        }}
        </script>
    </body>
    </html>
    """,
    height=350 # 稍微增加高度以容納更多文字
)
st.divider()

# --- 初始化 Session State ---
if 'slides_data' not in st.session_state:
    st.session_state['slides_data'] = []
if 'status_report' not in st.session_state:
    st.session_state['status_report'] = []

# --- 輔助函數：遍歷 Word ---
def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    else:
        raise ValueError("只支援讀取整份 Document")
    for child in parent_elm.iterchildren():
        if child.tag.endswith('p'):
            yield Paragraph(child, parent)
        elif child.tag.endswith('tbl'):
            yield Table(child, parent)

# --- 函數：搜尋 PDF 截圖 ---
def extract_specific_figure_from_pdf(pdf_stream, target_fig_text):
    if not target_fig_text:
        return None, "Word 中未指定代表圖文字"
    try:
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        pattern = re.compile(r'((?:FIG\.?|Figure|圖)\s*[0-9]+[A-Za-z]*)', re.IGNORECASE)
        search_keywords = []
        lines = target_fig_text.split('\n')
        for line in lines:
            match = pattern.search(line)
            if match:
                clean_keyword = match.group(1).replace(" ", "").upper()
                search_keywords.append(clean_keyword)
        
        if not search_keywords:
             first_line = lines[0].strip()
             if first_line:
                 search_keywords.append(first_line[:10].replace(" ", "").upper())

        target_keyword = search_keywords[0] if search_keywords else ""
        if not target_keyword:
            return None, "無法從說明文字中識別出圖號"

        found_page_index = None
        matched_keyword_log = ""
        for i, page in enumerate(doc):
            page_text = page.get_text().replace(" ", "").upper()
            if target_keyword in page_text:
                found_page_index = i
                matched_keyword_log = target_keyword
                break
        
        if found_page_index is not None:
            page = doc[found_page_index]
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            return pix.tobytes("png"), f"成功"
        return None, f"PDF 中找不到關鍵字「{target_keyword}」"
    except Exception as e:
        return None, f"PDF 解析發生錯誤: {str(e)}"

# --- 函數：提取專利號 ---
def extract_patent_number_from_text(text):
    clean_text = text.replace("：", ":").replace(" ", "")
    match = re.search(r'([a-zA-Z]{2,4}\d+[a-zA-Z]?)', clean_text)
    if match: return match.group(1)
    return ""

# --- 函數：提取日期 (排序用) ---
def extract_date_for_sort(text):
    match = re.search(r'(\d{4})[./-](\d{1,2})[./-](\d{1,2})', text)
    if match: return f"{match.group(1)}{match.group(2).zfill(2)}{match.group(3).zfill(2)}"
    return "99999999"

# --- 函數：提取公司 (排序用) ---
def extract_company_for_sort(text):
    lines = text.split('\n')
    for line in lines:
        if "公司" in line or "申請人" in line:
            if "案號" in line and "日期" in line: continue
            return line.replace("公司", "").replace("申請人", "").replace("：", "").replace(":", "").strip()
    return "ZZZ"

# --- 函數：解析 Word 檔案 (包含第6點解析) ---
def parse_word_file(uploaded_docx):
    try:
        doc = docx.Document(uploaded_docx)
        cases = []
        # 新增 claim_text 欄位
        current_case = {
            "case_info": "", "problem": "", "spirit": "", "key_point": "", "rep_fig_text": "", "claim_text": "",
            "image_data": None, "image_name": "Word匯入", "raw_case_no": "",
            "sort_date": "99999999", "sort_company": "ZZZ",
            "source_file": uploaded_docx.name, "missing_fields": []
        }
        current_field = None 
        
        all_lines = []
        for block in iter_block_items(doc):
            if isinstance(block, Paragraph):
                if block.text.strip(): all_lines.append(block.text.strip())
            elif isinstance(block, Table):
                for row in block.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            if p.text.strip(): all_lines.append(p.text.strip())
        
        for text in all_lines:
            # 1. 案號 (新案件起點)
            if "案號" in text or "索號" in text:
                if current_case["case_info"] and current_field != "case_info_block":
                    if not current_case["problem"]: current_case["missing_fields"].append("解決問題")
                    cases.append(current_case)
                    current_case = {
                        "case_info": "", "problem": "", "spirit": "", "key_point": "", "rep_fig_text": "", "claim_text": "",
                        "image_data": None, "image_name": "Word匯入", "raw_case_no": "",
                        "sort_date": "99999999", "sort_company": "ZZZ",
                        "source_file": uploaded_docx.name, "missing_fields": []
                    }
                current_field = "case_info_block"
                current_case["case_info"] = text
                extracted_no = extract_patent_number_from_text(text)
                if extracted_no: current_case["raw_case_no"] = extracted_no
                current_case["sort_date"] = extract_date_for_sort(text)
                current_case["sort_company"] = extract_company_for_sort(text)
                continue

            # 2. 欄位切換
            if "解決問題" in text:
                current_field = "problem"
                current_case["problem"] = re.sub(r'^[0-9.．]*\s*解決問題[:：]?\s*', '', text)
                continue
            elif "發明精神" in text:
                current_field = "spirit"
                current_case["spirit"] = re.sub(r'^[0-9.．]*\s*發明精神[:：]?\s*', '', text)
                continue
            elif "重點" in text:
                current_field = "key_point"
                current_case["key_point"] = re.sub(r'^[0-9.．]*\s*(一句)?重點[:：]?\s*', '', text)
                continue
            elif "代表圖" in text:
                current_field = "rep_fig"
                current_case["rep_fig_text"] = re.sub(r'^[0-9.．]*\s*代表圖[:：]?\s*', '', text).strip()
                continue
            # 新增：Claim 欄位辨識 (相容 "6.獨立項claim")
            elif "獨立項" in text or ("claim" in text.lower() and "6" in text):
                current_field = "claim"
                content = re.sub(r'^[0-9.．]*\s*(獨立項)?(claim)?[:：]?\s*', '', text, flags=re.IGNORECASE).strip()
                current_case["claim_text"] = content
                continue

            # 3. 內容填充
            if current_field == "case_info_block":
                current_case["case_info"] += "\n" + text
                if current_case["sort_date"] == "99999999": current_case["sort_date"] = extract_date_for_sort(text)
                extracted_comp = extract_company_for_sort(current_case["case_info"])
                if extracted_comp != "ZZZ": current_case["sort_company"] = extracted_comp
                if not current_case["raw_case_no"]:
                    extracted_no = extract_patent_number_from_text(text)
                    if extracted_no: current_case["raw_case_no"] = extracted_no
            elif current_field == "rep_fig":
                current_case["rep_fig_text"] += "\n" + text
            elif current_field == "problem":
                current_case["problem"] += "\n" + text
            elif current_field == "spirit":
                current_case["spirit"] += "\n" + text
            elif current_field == "key_point":
                current_case["key_point"] += "\n" + text
            elif current_field == "claim": 
                current_case["claim_text"] += "\n" + text

        if current_case["case_info"]:
            if not current_case["problem"]: current_case["missing_fields"].append("解決問題")
            cases.append(current_case)
        return cases
    except Exception as e:
        st.error(f"解析 Word 錯誤 ({uploaded_docx.name}): {e}")
        return []

# --- 側邊欄 ---
with st.sidebar:
    st.header("1. 匯入資料")
    word_files = st.file_uploader("Word 檔案 (可多選)", type=['docx'], accept_multiple_files=True)
    pdf_files = st.file_uploader("PDF 檔案 (可多選)", type=['pdf'], accept_multiple_files=True)
    
    st.divider()
    st.header("2. 輸出設定")
    add_claim_slide = st.checkbox("✅ 是否要產生 Claim 分頁", value=False, help="勾選後，每個案子會多出一頁專門放獨立項 Claim")

    if word_files and st.button("🔄 開始智能整合", type="primary"):
        all_cases = []
        status_report_list = []
        
        for wf in word_files:
            all_cases.extend(parse_word_file(wf))
        
        pdf_file_map = {}
        if pdf_files:
            for pf in pdf_files:
                clean = re.sub(r'[^a-zA-Z0-9]', '', pf.name.rsplit('.', 1)[0])
                pdf_file_map[clean] = pf.read()

        match_count = 0
        with st.spinner("處理中..."):
            for case in all_cases:
                case_key = case["raw_case_no"]
                target_fig = case["rep_fig_text"]
                status = {
                    "來源": case["source_file"], "案號": case_key if case_key else "?",
                    "公司": case["sort_company"], "日期": case["sort_date"],
                    "狀態": "未處理", "原因": "", "缺漏": ", ".join(case["missing_fields"])
                }
                
                matched_pdf = None
                for pk, pb in pdf_file_map.items():
                    if case_key and ((pk.lower() in case_key.lower()) or (case_key.lower() in pk.lower())):
                        if len(case_key) > 4: matched_pdf = pb; break
                
                if matched_pdf:
                    img_data, msg = extract_specific_figure_from_pdf(matched_pdf, target_fig)
                    if img_data:
                        case["image_data"] = img_data
                        status["狀態"] = "✅ 成功"; match_count += 1
                    else:
                        status["狀態"] = "⚠️ 缺圖"; status["原因"] = msg
                else:
                    if not target_fig: status["狀態"] = "⚠️ 缺資訊"; status["原因"] = "Word無代表圖"
                    else: status["狀態"] = "❌ 無PDF"; status["原因"] = f"找不到PDF: {case_key}"
                status_report_list.append(status)

        all_cases.sort(key=lambda x: (x["sort_company"].upper(), x["sort_date"]))
        status_report_list.sort(key=lambda x: (x["公司"].upper(), x["日期"]))

        if all_cases:
            st.session_state['slides_data'] = all_cases
            st.session_state['status_report'] = status_report_list
            st.success(f"完成！共 {len(all_cases)} 筆資料。")
        else:
            st.warning("無資料。")

    if st.session_state['slides_data']:
        st.divider()
        if st.button("🗑️ 清除重來"):
            st.session_state['slides_data'] = []
            st.session_state['status_report'] = []
            st.rerun()

# --- 主畫面 ---
if not st.session_state['slides_data']:
    st.info("👈 請先上傳檔案。")
else:
    st.subheader(f"📋 預覽 (已排序: 申請人 -> 日期)")
    cols = st.columns(3)
    for i, data in enumerate(st.session_state['slides_data']):
        with cols[i % 3]:
            with st.container(border=True):
                st.markdown(f"**Case {i+1}**")
                st.caption(f"{data['sort_company']} | {data['sort_date']}")
                st.text(data['case_info'][:80] + "...")
                if data['image_data']: st.image(data['image_data'], use_column_width=True)
                else: st.warning("無圖片")
                
                claim_preview = data['claim_text'][:50] + "..." if data['claim_text'] else "(無 Claim 資料)"
                st.caption(f"Claim: {claim_preview}")

    # --- PPT 生成邏輯 ---
    def generate_ppt(slides_data, need_claim_slide):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for data in slides_data:
            # === 第一頁 ===
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 左上
            left, top, width, height = Inches(0.5), Inches(0.5), Inches(5.0), Inches(2.0)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame; tf.word_wrap = True
            for line in data['case_info'].split('\n'):
                if line.strip():
                    p = tf.add_paragraph(); p.text = line.strip(); p.font.size = Pt(20); p.font.bold = True

            # 右上
            img_left = Inches(5.5); img_top = Inches(0.5); img_height = Inches(4.0); img_width = Inches(7.0)
            if data['image_data']:
                slide.shapes.add_picture(BytesIO(data['image_data']), img_left, img_top, height=img_height)
            else:
                txBox = slide.shapes.add_textbox(img_left, img_top, img_width, img_height)
                tf = txBox.text_frame; tf.word_wrap = True
                content = data['rep_fig_text'] if data['rep_fig_text'].strip() else "無代表圖資訊"
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph(); p.text = line.strip(); p.font.size = Pt(16)

            # 中下 & 底部
            left, top, width, height = Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame; tf.word_wrap = True
            p1 = tf.add_paragraph(); p1.text = "• 解決問題：" + data['problem']; p1.font.size = Pt(18); p1.space_after = Pt(12)
            p2 = tf.add_paragraph(); p2.text = "• 發明精神：" + data['spirit']; p2.font.size = Pt(18)

            left, top, width, height = Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 192, 0); shape.line.color.rgb = RGBColor(255, 192, 0)
            p = shape.text_frame.paragraphs[0]; p.text = data['key_point']; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(20); p.font.bold = True
            shape.text_frame.vertical_anchor = MSO_SHAPE.RECTANGLE

            # === 第二頁：Claim (勾選時) ===
            if need_claim_slide:
                slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 左上 (同上)
                left, top, width, height = Inches(0.5), Inches(0.5), Inches(5.0), Inches(2.0)
                txBox = slide2.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame; tf.word_wrap = True
                for line in data['case_info'].split('\n'):
                    if line.strip():
                        p = tf.add_paragraph(); p.text = line.strip(); p.font.size = Pt(20); p.font.bold = True
                
                # 中間：Claim
                left, top, width, height = Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.5)
                txBox = slide2.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame; tf.word_wrap = True
                
                p_title = tf.add_paragraph()
                p_title.text = "【獨立項 Claim】"
                p_title.font.size = Pt(24); p_title.font.bold = True; p_title.font.color.rgb = RGBColor(0, 112, 192)
                p_title.space_after = Pt(10)
                
                claim_content = data['claim_text'] if data['claim_text'].strip() else "(Word 中無 Claim 資料)"
                for line in claim_content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph(); p.text = line.strip(); p.font.size = Pt(18); p.space_after = Pt(6)

        return prs

    st.divider()
    if st.button("🚀 生成 PowerPoint (.pptx)", type="primary"):
        prs = generate_ppt(st.session_state['slides_data'], add_claim_slide)
        binary_output = BytesIO()
        prs.save(binary_output)
        binary_output.seek(0)
        st.download_button("📥 下載 PPT", binary_output, "slides_with_claims.pptx")

    st.divider()
    st.subheader("📊 診斷報告")
    if st.session_state['status_report']:
        st.dataframe(pd.DataFrame(st.session_state['status_report']), hide_index=True)
